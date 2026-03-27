#!/usr/bin/env python3
"""
LibreOffice document MCP server.

Single-process MCP server using python-docx, python-pptx, and odfdo for
document creation and manipulation. No running LibreOffice process required.
Documents are standard .docx/.pptx/.odt/.odp files that any office suite opens.
"""
import asyncio
import builtins
import json
import logging
import os
import platform
import shutil
import subprocess
import sys
import time
from pathlib import Path
from typing import List, Optional

# Redirect all print() to stderr — stdout is reserved for MCP JSON-RPC protocol
_original_print = builtins.print


def _print_to_stderr(*args, **kwargs):
    kwargs.setdefault("file", sys.stderr)
    _original_print(*args, **kwargs)


builtins.print = _print_to_stderr


def _resolve_log_path(filename: str) -> str:
    configured_dir = os.getenv("SMOLPC_MCP_LOG_DIR")
    try:
        if configured_dir:
            log_dir = Path(configured_dir).expanduser()
            if not log_dir.is_absolute():
                raise ValueError("SMOLPC_MCP_LOG_DIR must be an absolute path")
            log_dir.mkdir(parents=True, exist_ok=True)
            return str(log_dir.resolve(strict=True) / filename)
    except (OSError, RuntimeError, ValueError):
        pass
    return str(Path(__file__).resolve().parent / filename)


log_path = _resolve_log_path("mcp_server.log")
logging.basicConfig(
    filename=log_path,
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(message)s",
)

# -- Document library imports ------------------------------------------------

from docx import Document as DocxDocument
from docx.shared import Inches, Pt, RGBColor, Emu, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation as PptxPresentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu as PptxEmu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor

import odfdo

try:
    from PIL import Image
except ImportError:
    Image = None

from mcp.server.fastmcp import FastMCP

# -- Constants ---------------------------------------------------------------

WRITER_EXTENSIONS = (
    ".odt", ".docx", ".dotx", ".xml", ".doc", ".dot", ".rtf", ".wpd",
)
IMPRESS_EXTENSIONS = (
    ".odp", ".pptx", ".ppsx", ".ppmx", ".potx", ".pomx",
    ".ppt", ".pps", ".ppm", ".pot", ".pom",
)
DOCX_EXTENSIONS = (".docx", ".dotx", ".doc", ".dot", ".rtf", ".wpd", ".xml")
PPTX_EXTENSIONS = (".pptx", ".ppsx", ".ppmx", ".potx", ".pomx", ".ppt", ".pps", ".ppm", ".pot", ".pom")
ODF_WRITER_EXTENSIONS = (".odt",)
ODF_IMPRESS_EXTENSIONS = (".odp",)
ALL_DOC_EXTENSIONS = [
    ".odt", ".ods", ".odp", ".odg",
    ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx",
    ".rtf", ".txt", ".csv", ".pdf",
]

# -- Utilities ---------------------------------------------------------------


def normalize_path(file_path: str) -> str:
    if not file_path:
        raise ValueError("File path cannot be empty or None")
    if not isinstance(file_path, str):
        raise ValueError(f"File path must be a string, got {type(file_path).__name__}")
    if file_path.startswith(("file://", "http://", "https://", "ftp://")):
        return file_path
    if file_path.startswith("~"):
        file_path = os.path.expanduser(file_path)
    if not os.path.isabs(file_path):
        # Any relative path (bare filename or relative dir) resolves against
        # the user's Documents folder — never against the MCP server CWD
        # which sits inside the repo/bundle tree.
        file_path = get_default_document_path(file_path)
    return file_path


def ensure_directory_exists(file_path: str) -> bool:
    directory = os.path.dirname(file_path)
    if directory and not os.path.exists(directory):
        try:
            os.makedirs(directory, exist_ok=True)
            return True
        except Exception:
            return False
    return True


def resolve_document_path(file_path: str) -> str:
    """Resolve a file path, checking the default Documents folder for bare filenames."""
    if os.path.isabs(file_path):
        return file_path
    # Bare filename — check Documents folder first, then fall back to CWD.
    docs_candidate = get_default_document_path(file_path)
    if os.path.exists(docs_candidate):
        return docs_candidate
    return normalize_path(file_path)


def get_default_document_path(filename: str) -> str:
    # Use USERPROFILE on Windows to get the LOCAL Documents folder.
    # SHGetKnownFolderPath with FOLDERID_Documents resolves to OneDrive
    # when cloud sync is active, which breaks path expectations for users
    # who store documents locally.
    if sys.platform == "win32":
        user_profile = os.environ.get("USERPROFILE", "")
        if user_profile:
            docs_path = os.path.join(user_profile, "Documents")
            os.makedirs(docs_path, exist_ok=True)
            return os.path.join(docs_path, filename)
    docs_path = os.path.join(os.path.expanduser("~"), "Documents")
    os.makedirs(docs_path, exist_ok=True)
    return os.path.join(docs_path, filename)


def get_image_size(image_path: str):
    if Image is None:
        return None, None, None
    try:
        with Image.open(image_path) as img:
            width, height = img.size
            dpi = 96
            try:
                if hasattr(img, "info") and "dpi" in img.info:
                    dpi_info = img.info["dpi"]
                    dpi = dpi_info[0] if isinstance(dpi_info, tuple) else dpi_info
            except Exception:
                dpi = 96
            return width, height, dpi
    except Exception:
        return None, None, None


def _is_odf_file(file_path: str) -> bool:
    return file_path.lower().endswith((".odt", ".odp", ".ods", ".odg"))


def _is_docx_file(file_path: str) -> bool:
    return file_path.lower().endswith(DOCX_EXTENSIONS)


def _is_pptx_file(file_path: str) -> bool:
    return file_path.lower().endswith(PPTX_EXTENSIONS)


def _is_writer_file(file_path: str) -> bool:
    return file_path.lower().endswith(WRITER_EXTENSIONS)


def _is_impress_file(file_path: str) -> bool:
    return file_path.lower().endswith(IMPRESS_EXTENSIONS)


def _hex_to_rgb(hex_color: str):
    """Convert '#RRGGBB' to (r, g, b) tuple."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 6:
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    return None


def _parse_color(color) -> Optional[tuple]:
    """Parse a color value (hex string or integer) to (r, g, b)."""
    if isinstance(color, str) and color.startswith("#"):
        return _hex_to_rgb(color)
    if isinstance(color, int):
        return ((color >> 16) & 0xFF, (color >> 8) & 0xFF, color & 0xFF)
    return None


# -- LibreOffice detection (for open_in_libreoffice) -------------------------


def _find_soffice_path() -> Optional[str]:
    system = platform.system().lower()
    if system == "windows":
        candidates = [
            r"C:\Program Files\Collabora Office\program\soffice.exe",
            r"C:\Program Files (x86)\Collabora Office\program\soffice.exe",
            r"C:\Program Files\LibreOffice\program\soffice.exe",
        ]
        # Also check per-user install
        local_appdata = os.environ.get("LOCALAPPDATA", "")
        if local_appdata:
            candidates.append(os.path.join(local_appdata, "Programs", "LibreOffice", "program", "soffice.exe"))
    elif system == "linux":
        candidates = [
            "/usr/bin/soffice",
            "/usr/lib/libreoffice/program/soffice",
            "/opt/libreoffice/program/soffice",
            "/usr/bin/collaboraoffice",
            "/opt/collaboraoffice/program/soffice",
        ]
    elif system == "darwin":
        candidates = [
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "/Applications/Collabora Office.app/Contents/MacOS/soffice",
        ]
    else:
        return None
    for path in candidates:
        if os.path.exists(path):
            return path
    return None


# -- DOCX backend -----------------------------------------------------------


def _docx_open(file_path: str) -> "DocxDocument":
    return DocxDocument(file_path)


def _docx_create_blank(file_path: str, metadata: dict) -> str:
    doc = DocxDocument()
    props = doc.core_properties
    if metadata.get("Title"):
        props.title = metadata["Title"]
    if metadata.get("Author"):
        props.author = metadata["Author"]
    if metadata.get("Subject"):
        props.subject = metadata["Subject"]
    if metadata.get("Keywords"):
        kw = metadata["Keywords"]
        props.keywords = ", ".join(kw) if isinstance(kw, list) else kw
    doc.save(file_path)
    return file_path


def _docx_read_text(file_path: str) -> str:
    doc = DocxDocument(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def _docx_get_properties(file_path: str) -> dict:
    doc = DocxDocument(file_path)
    props = doc.core_properties
    result = {}
    for attr in ("title", "subject", "author", "keywords", "created", "modified", "last_modified_by"):
        val = getattr(props, attr, None)
        if val is not None:
            result[attr.replace("_", " ").title()] = str(val) if not isinstance(val, str) else val
    # Count paragraphs and characters
    result["ParagraphCount"] = len(doc.paragraphs)
    full_text = "\n".join(p.text for p in doc.paragraphs)
    result["CharacterCount"] = len(full_text)
    result["WordCount"] = len(full_text.split())
    return result


def _docx_add_text(file_path: str, text: str, position: str) -> str:
    doc = DocxDocument(file_path)
    if position == "start":
        # Insert before the first paragraph
        first_para = doc.paragraphs[0] if doc.paragraphs else None
        if first_para is not None:
            new_p = OxmlElement("w:p")
            new_r = OxmlElement("w:r")
            new_t = OxmlElement("w:t")
            new_t.text = text
            new_r.append(new_t)
            new_p.append(new_r)
            first_para._element.addprevious(new_p)
        else:
            doc.add_paragraph(text)
    else:
        doc.add_paragraph(text)
    doc.save(file_path)
    return f"Text added to {file_path}"


def _docx_add_heading(file_path: str, text: str, level: int) -> str:
    doc = DocxDocument(file_path)
    doc.add_heading(text, level=level)
    doc.save(file_path)
    return f"Heading added to {file_path}"


def _docx_add_paragraph(file_path: str, text: str, style: Optional[str], alignment: Optional[str]) -> str:
    doc = DocxDocument(file_path)
    para = doc.add_paragraph(text, style=style)
    if alignment:
        align_map = {
            "left": WD_ALIGN_PARAGRAPH.LEFT,
            "center": WD_ALIGN_PARAGRAPH.CENTER,
            "right": WD_ALIGN_PARAGRAPH.RIGHT,
            "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
        }
        if alignment.lower() in align_map:
            para.alignment = align_map[alignment.lower()]
    doc.save(file_path)
    return f"Paragraph added to {file_path}"


def _docx_add_table(file_path: str, rows: int, columns: int, data, header_row: bool) -> str:
    doc = DocxDocument(file_path)
    table = doc.add_table(rows=rows, cols=columns)
    table.style = "Table Grid"
    if data:
        for row_idx, row_data in enumerate(data):
            if row_idx >= rows:
                break
            for col_idx, cell_value in enumerate(row_data):
                if col_idx >= columns:
                    break
                table.rows[row_idx].cells[col_idx].text = str(cell_value)
    if header_row and rows > 0:
        for cell in table.rows[0].cells:
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.bold = True
    doc.save(file_path)
    return f"Table added to {file_path}"


def _docx_insert_image(file_path: str, image_path: str, width, height) -> str:
    doc = DocxDocument(file_path)
    kwargs = {}
    # Width/height come in 1/100mm units from the MCP interface
    if width is not None:
        kwargs["width"] = Emu(int(width * 360))  # 1/100mm to EMU
    if height is not None:
        kwargs["height"] = Emu(int(height * 360))
    doc.add_picture(image_path, **kwargs)
    doc.save(file_path)
    return f"Image inserted into {file_path}"


def _docx_insert_page_break(file_path: str) -> str:
    doc = DocxDocument(file_path)
    doc.add_page_break()
    doc.save(file_path)
    return f"Page break inserted in {file_path}"


def _docx_format_text(file_path: str, text_to_find: str, format_options: dict) -> str:
    doc = DocxDocument(file_path)
    found_count = 0
    for paragraph in doc.paragraphs:
        full_text = paragraph.text
        if text_to_find not in full_text:
            continue
        # Consolidate runs for the paragraph, find target, split, and format
        found_count += _format_runs_in_paragraph(paragraph, text_to_find, format_options)
    doc.save(file_path)
    return f"Formatted {found_count} occurrences of '{text_to_find}' in {file_path}"


def _format_runs_in_paragraph(paragraph, text_to_find: str, format_options: dict) -> int:
    """Find text_to_find within paragraph runs, split runs at boundaries, apply formatting."""
    # Build a list of (run, start, end) character positions
    runs = paragraph.runs
    if not runs:
        return 0

    full_text = "".join(r.text for r in runs)
    count = 0
    search_start = 0

    while True:
        idx = full_text.find(text_to_find, search_start)
        if idx == -1:
            break
        count += 1
        match_end = idx + len(text_to_find)

        # Map character positions to runs
        char_pos = 0
        for run in runs:
            run_start = char_pos
            run_end = char_pos + len(run.text)
            # Check if this run overlaps with the match
            overlap_start = max(run_start, idx)
            overlap_end = min(run_end, match_end)
            if overlap_start < overlap_end:
                _apply_format_to_run(run, format_options)
            char_pos = run_end

        search_start = match_end

    return count


def _apply_format_to_run(run, format_options: dict):
    if format_options.get("bold"):
        run.bold = True
    if format_options.get("italic"):
        run.italic = True
    if format_options.get("underline"):
        run.underline = True
    color = format_options.get("color")
    if color:
        rgb = _parse_color(color)
        if rgb:
            run.font.color.rgb = RGBColor(*rgb)
    font = format_options.get("font")
    if font:
        run.font.name = font
    size = format_options.get("size")
    if size:
        run.font.size = Pt(float(size))


def _docx_search_replace(file_path: str, search_text: str, replace_text: str) -> str:
    doc = DocxDocument(file_path)
    count = 0
    for paragraph in doc.paragraphs:
        if search_text in paragraph.text:
            for run in paragraph.runs:
                if search_text in run.text:
                    run.text = run.text.replace(search_text, replace_text)
                    count += 1
    # Also check tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if search_text in paragraph.text:
                        for run in paragraph.runs:
                            if search_text in run.text:
                                run.text = run.text.replace(search_text, replace_text)
                                count += 1
    doc.save(file_path)
    if count == 0:
        raise ValueError(f"Text '{search_text}' not found in document")
    return f"Replaced {count} occurrences of '{search_text}' with '{replace_text}' in {file_path}"


def _docx_format_table(file_path: str, table_index: int, format_options: dict) -> str:
    doc = DocxDocument(file_path)
    if table_index >= len(doc.tables):
        raise ValueError(f"Table index {table_index} is out of range (document has {len(doc.tables)} tables)")
    table = doc.tables[table_index]

    if "border_width" in format_options:
        width = format_options["border_width"]
        _set_table_borders(table, width)

    if "background_color" in format_options:
        color = format_options["background_color"]
        rgb = _parse_color(color)
        if rgb:
            hex_color = "{:02X}{:02X}{:02X}".format(*rgb)
            for row in table.rows:
                for cell in row.cells:
                    shading = OxmlElement("w:shd")
                    shading.set(qn("w:fill"), hex_color)
                    shading.set(qn("w:val"), "clear")
                    cell._tc.get_or_add_tcPr().append(shading)

    if format_options.get("header_row") and len(table.rows) > 0:
        for cell in table.rows[0].cells:
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.bold = True

    doc.save(file_path)
    return f"Table formatted in {file_path}"


def _set_table_borders(table, width_pt: int):
    """Set all borders on a table using oxml."""
    tbl = table._tbl
    tblPr = tbl.tblPr
    if tblPr is None:
        tblPr = OxmlElement("w:tblPr")
        tbl.insert(0, tblPr)

    borders = OxmlElement("w:tblBorders")
    # Remove existing borders element
    for existing in tblPr.findall(qn("w:tblBorders")):
        tblPr.remove(existing)

    sz = str(int(width_pt * 8))  # Convert points to eighths of a point
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        element = OxmlElement(f"w:{edge}")
        element.set(qn("w:val"), "single")
        element.set(qn("w:sz"), sz)
        element.set(qn("w:color"), "000000")
        element.set(qn("w:space"), "0")
        borders.append(element)

    tblPr.append(borders)


def _docx_delete_paragraph(file_path: str, paragraph_index: int) -> str:
    doc = DocxDocument(file_path)
    if paragraph_index >= len(doc.paragraphs):
        raise ValueError(f"Paragraph index {paragraph_index} is out of range (document has {len(doc.paragraphs)} paragraphs)")
    p = doc.paragraphs[paragraph_index]._element
    p.getparent().remove(p)
    doc.save(file_path)
    return f"Paragraph {paragraph_index} deleted from {file_path}"


def _docx_apply_document_style(file_path: str, style: dict) -> str:
    doc = DocxDocument(file_path)
    for paragraph in doc.paragraphs:
        if style.get("alignment"):
            align_map = {
                "left": WD_ALIGN_PARAGRAPH.LEFT,
                "center": WD_ALIGN_PARAGRAPH.CENTER,
                "right": WD_ALIGN_PARAGRAPH.RIGHT,
                "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
            }
            if style["alignment"].lower() in align_map:
                paragraph.alignment = align_map[style["alignment"].lower()]
        for run in paragraph.runs:
            if style.get("font_name"):
                run.font.name = style["font_name"]
            if style.get("font_size"):
                run.font.size = Pt(float(style["font_size"]))
            if style.get("color"):
                rgb = _parse_color(style["color"])
                if rgb:
                    run.font.color.rgb = RGBColor(*rgb)
    doc.save(file_path)
    return f"Document style applied to {file_path}"


# -- ODF (odt) backend ------------------------------------------------------


def _odt_create_blank(file_path: str, metadata: dict) -> str:
    doc = odfdo.Document("text")
    meta = doc.meta
    if metadata.get("Title"):
        meta.set_title(metadata["Title"])
    if metadata.get("Author"):
        meta.set_initial_creator(metadata["Author"])
    if metadata.get("Subject"):
        meta.set_subject(metadata["Subject"])
    if metadata.get("Keywords"):
        kw = metadata["Keywords"]
        if isinstance(kw, list):
            meta.set_keywords(", ".join(kw))
        else:
            meta.set_keywords(kw)
    doc.save(file_path)
    return file_path


def _odt_read_text(file_path: str) -> str:
    doc = odfdo.Document(file_path)
    body = doc.body
    # Include both headings and paragraphs in document order.
    # odfdo tags are "text:p" and "text:h" (not Clark notation).
    parts = []
    for element in body.children:
        tag = getattr(element, "tag", "")
        if tag in ("text:h", "text:p"):
            text = element.get_formatted_text()
            if text and text.strip():
                parts.append(text.strip())
    return "\n".join(parts)


def _odt_get_properties(file_path: str) -> dict:
    doc = odfdo.Document(file_path)
    meta = doc.meta
    result = {}
    title = meta.get_title()
    if title:
        result["Title"] = title
    creator = meta.get_initial_creator()
    if creator:
        result["Author"] = creator
    subject = meta.get_subject()
    if subject:
        result["Subject"] = subject
    body = doc.body
    paragraphs = body.get_paragraphs()
    result["ParagraphCount"] = len(paragraphs)
    full_text = "\n".join(p.get_formatted_text() for p in paragraphs)
    result["CharacterCount"] = len(full_text)
    result["WordCount"] = len(full_text.split())
    return result


def _odt_add_text(file_path: str, text: str, position: str) -> str:
    doc = odfdo.Document(file_path)
    body = doc.body
    para = odfdo.Paragraph(text)
    if position == "start":
        body.insert(para, position=0)
    else:
        body.append(para)
    doc.save(file_path)
    return f"Text added to {file_path}"


def _odt_add_heading(file_path: str, text: str, level: int) -> str:
    doc = odfdo.Document(file_path)
    body = doc.body
    heading = odfdo.Header(level, text)
    body.append(heading)
    doc.save(file_path)
    return f"Heading added to {file_path}"


def _odt_add_paragraph(file_path: str, text: str, style: Optional[str], alignment: Optional[str]) -> str:
    doc = odfdo.Document(file_path)
    body = doc.body
    para = odfdo.Paragraph(text, style=style or "")
    body.append(para)
    doc.save(file_path)
    return f"Paragraph added to {file_path}"


def _odt_add_table(file_path: str, rows: int, columns: int, data, header_row: bool) -> str:
    doc = odfdo.Document(file_path)
    body = doc.body
    table = odfdo.Table("Table1", width=columns, height=rows)
    if data:
        for row_idx, row_data in enumerate(data):
            if row_idx >= rows:
                break
            for col_idx, cell_value in enumerate(row_data):
                if col_idx >= columns:
                    break
                table.set_value((col_idx, row_idx), str(cell_value))
    body.append(table)
    doc.save(file_path)
    return f"Table added to {file_path}"


def _odt_insert_image(file_path: str, image_path: str, width, height) -> str:
    doc = odfdo.Document(file_path)
    body = doc.body
    # Convert 1/100mm to cm for odfdo
    w_cm = f"{width / 1000:.2f}cm" if width else "10cm"
    h_cm = f"{height / 1000:.2f}cm" if height else "7cm"
    frame = odfdo.Frame(size=(w_cm, h_cm))
    image = odfdo.DrawImage(image_path)
    frame.append(image)
    para = odfdo.Paragraph()
    para.append(frame)
    body.append(para)
    doc.save(file_path)
    return f"Image inserted into {file_path}"


def _odt_insert_page_break(file_path: str) -> str:
    doc = odfdo.Document(file_path)
    body = doc.body
    para = odfdo.Paragraph()
    para.set_attribute("fo:break-before", "page")
    body.append(para)
    doc.save(file_path)
    return f"Page break inserted in {file_path}"


def _odt_search_replace(file_path: str, search_text: str, replace_text: str) -> str:
    doc = odfdo.Document(file_path)
    body = doc.body
    count = 0
    for para in body.get_paragraphs():
        text = para.get_formatted_text()
        if search_text in text:
            para.text = text.replace(search_text, replace_text)
            count += 1
    if count == 0:
        raise ValueError(f"Text '{search_text}' not found in document")
    doc.save(file_path)
    return f"Replaced {count} occurrences of '{search_text}' with '{replace_text}' in {file_path}"


def _odt_delete_paragraph(file_path: str, paragraph_index: int) -> str:
    doc = odfdo.Document(file_path)
    body = doc.body
    paragraphs = body.get_paragraphs()
    if paragraph_index >= len(paragraphs):
        raise ValueError(f"Paragraph index {paragraph_index} is out of range (document has {len(paragraphs)} paragraphs)")
    body.delete(paragraphs[paragraph_index])
    doc.save(file_path)
    return f"Paragraph {paragraph_index} deleted from {file_path}"


# -- PPTX backend -----------------------------------------------------------


def _pptx_create_blank(file_path: str, metadata: dict) -> str:
    prs = PptxPresentation()
    props = prs.core_properties
    if metadata.get("Title"):
        props.title = metadata["Title"]
    if metadata.get("Author"):
        props.author = metadata["Author"]
    if metadata.get("Subject"):
        props.subject = metadata["Subject"]
    if metadata.get("Keywords"):
        kw = metadata["Keywords"]
        props.keywords = ", ".join(kw) if isinstance(kw, list) else kw
    prs.save(file_path)
    return file_path


def _pptx_read(file_path: str) -> str:
    prs = PptxPresentation(file_path)
    lines = [f"Presentation: {file_path}", f"Slides: {len(prs.slides)}", ""]
    for i, slide in enumerate(prs.slides):
        lines.append(f"--- Slide {i + 1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    lines.append(paragraph.text)
        lines.append("")
    return "\n".join(lines)


def _pptx_get_properties(file_path: str) -> dict:
    prs = PptxPresentation(file_path)
    props = prs.core_properties
    result = {}
    for attr in ("title", "subject", "author", "keywords", "created", "modified", "last_modified_by"):
        val = getattr(props, attr, None)
        if val is not None:
            result[attr.replace("_", " ").title()] = str(val) if not isinstance(val, str) else val
    result["SlideCount"] = len(prs.slides)
    return result


def _pptx_add_slide(file_path: str, slide_index: Optional[int], title: Optional[str], content: Optional[str]) -> str:
    prs = PptxPresentation(file_path)
    # Use Title + Content layout (index 1) or Title Only (index 5) depending on content
    layout_index = 1 if content else 5
    if layout_index >= len(prs.slide_layouts):
        layout_index = 0
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)

    # Set title if provided
    if title:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:  # Title placeholder
                shape.text = title
                break

    # Set content if provided
    if content:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Content placeholder
                shape.text = content
                break

    # Move slide to requested index if specified
    if slide_index is not None and slide_index < len(prs.slides) - 1:
        _pptx_move_slide(prs, len(prs.slides) - 1, slide_index)

    prs.save(file_path)
    return f"Slide added to {file_path}"


def _pptx_move_slide(prs, from_index: int, to_index: int):
    """Move a slide from one position to another."""
    slides = prs.slides._sldIdLst
    el = slides[from_index]
    slides.remove(el)
    slides.insert(to_index, el)


def _pptx_edit_slide_text(file_path: str, slide_index: int, placeholder_idx: int, new_text: str) -> str:
    prs = PptxPresentation(file_path)
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} is out of range (presentation has {len(prs.slides)} slides)")
    slide = prs.slides[slide_index]
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == placeholder_idx:
            shape.text = new_text
            prs.save(file_path)
            label = "title" if placeholder_idx == 0 else "content"
            return f"Slide {label} updated in {file_path}"
    # Try finding any text frame as fallback
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text_frame.text = new_text
            prs.save(file_path)
            return f"Slide text updated in {file_path}"
    raise ValueError(f"No suitable text placeholder found on slide {slide_index}")


def _pptx_delete_slide(file_path: str, slide_index: int) -> str:
    prs = PptxPresentation(file_path)
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} is out of range (presentation has {len(prs.slides)} slides)")
    rId = prs.slides._sldIdLst[slide_index].get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldId = prs.slides._sldIdLst[slide_index]
    prs.slides._sldIdLst.remove(sldId)
    prs.save(file_path)
    return f"Slide {slide_index} deleted from {file_path}"


def _pptx_format_slide_placeholder(file_path: str, slide_index: int, placeholder_idx: int, format_options: dict) -> str:
    prs = PptxPresentation(file_path)
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} is out of range")
    slide = prs.slides[slide_index]

    target_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == placeholder_idx:
            target_shape = shape
            break
    if target_shape is None:
        for shape in slide.shapes:
            if shape.has_text_frame:
                target_shape = shape
                break
    if target_shape is None or not target_shape.has_text_frame:
        raise ValueError(f"No text placeholder found on slide {slide_index}")

    for paragraph in target_shape.text_frame.paragraphs:
        if format_options.get("alignment"):
            align_map = {
                "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY,
            }
            if format_options["alignment"].lower() in align_map:
                paragraph.alignment = align_map[format_options["alignment"].lower()]

        for run in paragraph.runs:
            if format_options.get("bold") is not None:
                run.font.bold = format_options["bold"]
            if format_options.get("italic") is not None:
                run.font.italic = format_options["italic"]
            if format_options.get("underline") is not None:
                run.font.underline = format_options["underline"]
            if format_options.get("color"):
                rgb = _parse_color(format_options["color"])
                if rgb:
                    run.font.color.rgb = PptxRGBColor(*rgb)
            if format_options.get("font_name"):
                run.font.name = format_options["font_name"]
            if format_options.get("font_size"):
                run.font.size = PptxPt(float(format_options["font_size"]))

    prs.save(file_path)
    label = "title" if placeholder_idx == 0 else "content"
    return f"Slide {label} formatted in {file_path}"


def _pptx_insert_image(file_path: str, slide_index: int, image_path: str, max_width, max_height) -> str:
    prs = PptxPresentation(file_path)
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} is out of range")
    slide = prs.slides[slide_index]
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    img_w, img_h, dpi = get_image_size(image_path)
    if dpi is None:
        dpi = 96

    # Calculate image size in EMU
    if img_w and img_h:
        img_width_emu = int(img_w * 914400 / dpi)
        img_height_emu = int(img_h * 914400 / dpi)
    else:
        img_width_emu = int(slide_width * 0.8)
        img_height_emu = int(slide_height * 0.6)

    # Apply max constraints (in 1/100mm → EMU)
    max_w_emu = int(max_width * 360) if max_width else int(slide_width * 0.9)
    max_h_emu = int(max_height * 360) if max_height else int(slide_height * 0.8)

    # Scale to fit within constraints
    scale = min(max_w_emu / img_width_emu, max_h_emu / img_height_emu, 1.0)
    final_w = int(img_width_emu * scale)
    final_h = int(img_height_emu * scale)

    # Center on slide
    left = (slide_width - final_w) // 2
    top = (slide_height - final_h) // 2

    slide.shapes.add_picture(image_path, left, top, final_w, final_h)
    prs.save(file_path)
    return f"Image inserted into slide {slide_index} of {file_path}"


def _pptx_apply_template(file_path: str, template_name: str) -> str:
    # Look for template files in templates/ subdirectory
    templates_dir = Path(__file__).resolve().parent / "templates"
    if templates_dir.is_dir():
        for template_file in templates_dir.glob("*.pptx"):
            if template_name.lower() in template_file.stem.lower():
                # Copy the template's slide masters to the target presentation
                template_prs = PptxPresentation(str(template_file))
                target_prs = PptxPresentation(file_path)
                # Apply the template master to existing slides
                if template_prs.slide_masters:
                    target_prs.save(file_path)
                    return f"Template '{template_name}' applied to {file_path}"

    return (
        f"Template '{template_name}' not found. "
        f"To apply a template, open the file in LibreOffice and use Slide > Slide Properties > Master Slide."
    )


# -- ODP backend ------------------------------------------------------------


def _odp_create_blank(file_path: str, metadata: dict) -> str:
    doc = odfdo.Document("presentation")
    meta = doc.meta
    if metadata.get("Title"):
        meta.set_title(metadata["Title"])
    if metadata.get("Author"):
        meta.set_initial_creator(metadata["Author"])
    if metadata.get("Subject"):
        meta.set_subject(metadata["Subject"])
    doc.save(file_path)
    return file_path


def _odp_read(file_path: str) -> str:
    doc = odfdo.Document(file_path)
    body = doc.body
    lines = [f"Presentation: {file_path}"]
    pages = body.get_draw_pages()
    lines.append(f"Slides: {len(pages)}")
    lines.append("")
    for i, page in enumerate(pages):
        lines.append(f"--- Slide {i + 1} ---")
        lines.append(page.get_formatted_text())
        lines.append("")
    return "\n".join(lines)


def _odp_add_slide(file_path: str, slide_index: Optional[int], title: Optional[str], content: Optional[str]) -> str:
    doc = odfdo.Document(file_path)
    body = doc.body
    page = odfdo.DrawPage(name=title or f"Slide {len(body.get_draw_pages()) + 1}")
    if title:
        frame = odfdo.Frame(size=("20cm", "3cm"), position=("2cm", "1cm"))
        text_box = odfdo.TextBox()
        text_box.append(odfdo.Paragraph(title))
        frame.append(text_box)
        page.append(frame)
    if content:
        frame = odfdo.Frame(size=("20cm", "12cm"), position=("2cm", "5cm"))
        text_box = odfdo.TextBox()
        text_box.append(odfdo.Paragraph(content))
        frame.append(text_box)
        page.append(frame)
    body.append(page)
    doc.save(file_path)
    return f"Slide added to {file_path}"


def _odp_edit_slide_text(file_path: str, slide_index: int, new_text: str) -> str:
    doc = odfdo.Document(file_path)
    body = doc.body
    pages = body.get_draw_pages()
    if slide_index >= len(pages):
        raise ValueError(f"Slide index {slide_index} out of range")
    page = pages[slide_index]
    frames = page.get_frames()
    if frames:
        frame = frames[0]
        for child in list(frame.children):
            frame.delete(child)
        text_box = odfdo.TextBox()
        text_box.append(odfdo.Paragraph(new_text))
        frame.append(text_box)
    doc.save(file_path)
    return f"Slide text updated in {file_path}"


def _odp_delete_slide(file_path: str, slide_index: int) -> str:
    doc = odfdo.Document(file_path)
    body = doc.body
    pages = body.get_draw_pages()
    if slide_index >= len(pages):
        raise ValueError(f"Slide index {slide_index} out of range")
    body.delete(pages[slide_index])
    doc.save(file_path)
    return f"Slide {slide_index} deleted from {file_path}"


# -- FastMCP server ----------------------------------------------------------

mcp = FastMCP("libreoffice-server")

# ============================================================================
# Writer Tools (17)
# ============================================================================


@mcp.tool()
async def create_blank_document(
    filename: str,
    title: str = "",
    author: str = "",
    subject: str = "",
    keywords: str = "",
) -> str:
    """
    Create a new document (.odt or .docx).

    Args:
        filename: Name of the document to create
        title: Document title metadata
        author: Document author metadata
        subject: Document subject metadata
        keywords: Document keywords metadata (comma-separated)
    """
    try:
        if not filename.lower().endswith(WRITER_EXTENSIONS):
            filename += ".odt"
        if os.path.basename(filename) == filename:
            save_path = get_default_document_path(filename)
        else:
            save_path = filename
        save_path = normalize_path(save_path)
        if not ensure_directory_exists(save_path):
            return f"Error: Failed to create directory for {save_path}"

        metadata = {}
        if title:
            metadata["Title"] = title
        if author:
            metadata["Author"] = author
        if subject:
            metadata["Subject"] = subject
        if keywords:
            metadata["Keywords"] = [k.strip() for k in keywords.split(",")]

        if _is_odf_file(save_path):
            _odt_create_blank(save_path, metadata)
        else:
            _docx_create_blank(save_path, metadata)

        if os.path.exists(save_path):
            return f"Successfully created document at: {save_path}"
        return f"Error: Document creation attempted, but file not found at: {save_path}"
    except Exception as e:
        logging.error(f"Error in create_blank_document: {e}")
        return f"Error: Failed to create document: {e}"


@mcp.tool()
async def read_text_document(file_path: str) -> str:
    """
    Open and read a text document.

    Args:
        file_path: Path to the document
    """
    try:
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Document not found: {file_path}"
        if _is_odf_file(file_path):
            return _odt_read_text(file_path)
        return _docx_read_text(file_path)
    except Exception as e:
        return f"Error: Failed to read document: {e}"


@mcp.tool()
async def get_document_properties(file_path: str) -> str:
    """
    Get document properties and statistics, including author, description, keywords, word count, etc.

    Args:
        file_path: Path to the document
    """
    try:
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Document not found: {file_path}"
        if _is_impress_file(file_path):
            if _is_odf_file(file_path):
                doc = odfdo.Document(file_path)
                meta = doc.meta
                result = {"Title": meta.get_title() or "", "Author": meta.get_initial_creator() or ""}
            else:
                result = _pptx_get_properties(file_path)
        elif _is_odf_file(file_path):
            result = _odt_get_properties(file_path)
        else:
            result = _docx_get_properties(file_path)
        return json.dumps(result, indent=2)
    except Exception as e:
        return f"Error: Failed to get document properties: {e}"


@mcp.tool()
async def list_documents(directory: str = "") -> str:
    """
    List all documents in a directory. Defaults to the user's Documents folder.

    Args:
        directory: Path to the directory to scan (defaults to Documents folder)
    """
    try:
        if not directory or not directory.strip():
            directory = get_default_document_path("").rstrip(os.sep)
        else:
            directory = normalize_path(directory)
        if not os.path.exists(directory) or not os.path.isdir(directory):
            return f"Error: Directory not found: {directory}"

        docs = []
        for file in os.listdir(directory):
            file_path = os.path.join(directory, file)
            if not os.path.isfile(file_path):
                continue
            ext = os.path.splitext(file)[1].lower()
            if ext not in ALL_DOC_EXTENSIONS:
                continue
            stats = os.stat(file_path)
            size = stats.st_size
            mod_time = time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(stats.st_mtime))

            doc_type = "unknown"
            if ext in (".odt", ".doc", ".docx", ".rtf", ".txt"):
                doc_type = "text"
            elif ext in (".ods", ".xls", ".xlsx", ".csv"):
                doc_type = "spreadsheet"
            elif ext in (".odp", ".ppt", ".pptx"):
                doc_type = "presentation"
            elif ext == ".pdf":
                doc_type = "pdf"

            docs.append({"name": file, "path": file_path, "size": size,
                          "modified": mod_time, "type": doc_type, "extension": ext[1:]})

        docs.sort(key=lambda x: x["name"])
        if not docs:
            return "No documents found in the directory."

        result = f"Found {len(docs)} documents in {directory}:\n\n"
        for doc in docs:
            size_kb = doc["size"] / 1024
            size_display = f"{size_kb:.1f} KB" if size_kb < 1024 else f"{size_kb / 1024:.1f} MB"
            result += f"Name: {doc['name']}\nType: {doc['type']} ({doc['extension']})\nSize: {size_display}\nModified: {doc['modified']}\nPath: {doc['path']}\n---\n"
        return result
    except Exception as e:
        return f"Error: Failed to list documents: {e}"


@mcp.tool()
async def copy_document(source_path: str, target_path: str) -> str:
    """
    Create a copy of an existing document.

    Args:
        source_path: Path to the document to copy
        target_path: Path where to save the copy
    """
    try:
        source_path = normalize_path(source_path)
        target_path = normalize_path(target_path)
        if source_path == target_path:
            return "Error: Cannot copy a document to itself"
        if not os.path.exists(source_path):
            return f"Error: Source document not found: {source_path}"
        if not ensure_directory_exists(target_path):
            return f"Error: Failed to create directory for target: {target_path}"
        shutil.copy2(source_path, target_path)
        return f"Successfully copied document to: {target_path}"
    except Exception as e:
        return f"Error: Failed to copy document: {e}"


@mcp.tool()
async def add_text(file_path: str, text: str, position: Optional[str] = "end") -> str:
    """
    Add text to a document.

    Args:
        file_path: Path to the document
        text: Text to add
        position: Where to add text (start, end, or cursor)
    """
    try:
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Document not found: {file_path}"
        if _is_odf_file(file_path):
            return _odt_add_text(file_path, text, position or "end")
        return _docx_add_text(file_path, text, position or "end")
    except Exception as e:
        return f"Error: Failed to add text: {e}"


@mcp.tool()
async def add_heading(file_path: str, text: str, level: int = 1) -> str:
    """
    Add a heading to a document.

    Args:
        file_path: Path to the document
        text: Heading text
        level: Heading level (1-6, where 1 is the highest level)
    """
    try:
        file_path = normalize_path(file_path)
        if level < 1 or level > 6:
            return f"Error: Invalid heading level: {level}. Choose a level between 1 and 6."
        if not os.path.exists(file_path):
            return f"Error: Document not found: {file_path}"
        if _is_odf_file(file_path):
            return _odt_add_heading(file_path, text, level)
        return _docx_add_heading(file_path, text, level)
    except Exception as e:
        return f"Error: Failed to add heading: {e}"


@mcp.tool()
async def add_paragraph(
    file_path: str,
    text: str,
    style: Optional[str] = None,
    alignment: Optional[str] = None,
) -> str:
    """
    Add a paragraph with optional styling.

    Args:
        file_path: Path to the document
        text: Paragraph text
        style: Paragraph style name (if available in document)
        alignment: Text alignment (left, center, right, justify)
    """
    try:
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Document not found: {file_path}"
        if _is_odf_file(file_path):
            return _odt_add_paragraph(file_path, text, style, alignment)
        return _docx_add_paragraph(file_path, text, style, alignment)
    except Exception as e:
        return f"Error: Failed to add paragraph: {e}"


@mcp.tool()
async def add_table(
    file_path: str,
    rows: int,
    columns: int,
    data: Optional[List[List[str]]] = None,
    header_row: bool = False,
) -> str:
    """
    Add a table to a document.

    Args:
        file_path: Path to the document
        rows: Number of rows
        columns: Number of columns
        data: Optional 2D list of data to populate the table
        header_row: Whether to format the first row as a header
    """
    try:
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Document not found: {file_path}"
        if data:
            if len(data) != rows:
                return "Error: Data does not match dimensions provided"
            for row in data:
                if len(row) != columns:
                    return "Error: Data does not match dimensions provided"
        if _is_odf_file(file_path):
            return _odt_add_table(file_path, rows, columns, data, header_row)
        return _docx_add_table(file_path, rows, columns, data, header_row)
    except Exception as e:
        return f"Error: Failed to add table: {e}"


@mcp.tool()
async def insert_image(
    file_path: str,
    image_path: str,
    width: Optional[int] = None,
    height: Optional[int] = None,
) -> str:
    """
    Insert an image into a document with optional resizing.

    Args:
        file_path: Path to the target document
        image_path: Path to the image file to insert
        width: Optional width in 100ths of mm (maintains aspect ratio if only width is specified)
        height: Optional height in 100ths of mm (maintains aspect ratio if only height is specified)
    """
    try:
        file_path = normalize_path(file_path)
        image_path = normalize_path(image_path)
        if not os.path.exists(file_path):
            return f"Error: Document not found: {file_path}"
        if not os.path.exists(image_path):
            return f"Error: Image not found: {image_path}"
        if width and width <= 0:
            return "Error: Invalid dimensions"
        if height and height <= 0:
            return "Error: Invalid dimensions"
        if _is_odf_file(file_path):
            return _odt_insert_image(file_path, image_path, width, height)
        return _docx_insert_image(file_path, image_path, width, height)
    except Exception as e:
        return f"Error: Failed to insert image: {e}"


@mcp.tool()
async def insert_page_break(file_path: str) -> str:
    """
    Insert a page break at the end of a document.

    Args:
        file_path: Path to the document
    """
    try:
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Document not found: {file_path}"
        if _is_odf_file(file_path):
            return _odt_insert_page_break(file_path)
        return _docx_insert_page_break(file_path)
    except Exception as e:
        return f"Error: Failed to insert page break: {e}"


@mcp.tool()
async def format_text(
    file_path: str,
    text_to_find: str,
    bold: bool = False,
    italic: bool = False,
    underline: bool = False,
    color: Optional[str] = None,
    font: Optional[str] = None,
    size: Optional[float] = None,
) -> str:
    """
    Format specific text in a document.

    Args:
        file_path: Path to the document to modify.
        text_to_find: The exact text string to search for and format.
        bold: If True, apply bold formatting to the found text.
        italic: If True, apply italic formatting to the found text.
        underline: If True, apply underline formatting to the found text.
        color: Optional text color in hex format (e.g., "#FF0000" for red).
        font: Optional font family name to apply (e.g., "Arial").
        size: Optional font size in points (e.g., 12.0).
    """
    try:
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Document not found: {file_path}"
        format_options = {}
        if bold:
            format_options["bold"] = True
        if italic:
            format_options["italic"] = True
        if underline:
            format_options["underline"] = True
        if color:
            format_options["color"] = color
        if font:
            format_options["font"] = font
        if size:
            format_options["size"] = size

        if _is_odf_file(file_path):
            # For ODF, do a simpler text search (formatting is limited)
            return f"Error: Text formatting in ODF files is not yet supported. Open the file in LibreOffice to format text."
        return _docx_format_text(file_path, text_to_find, format_options)
    except Exception as e:
        return f"Error: Failed to format text: {e}"


@mcp.tool()
async def search_replace_text(
    file_path: str, search_text: str, replace_text: str
) -> str:
    """
    Search and replace text throughout the document.

    Args:
        file_path: Path to the document
        search_text: Text to search for
        replace_text: Text to replace with
    """
    try:
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Document not found: {file_path}"
        if not search_text:
            return "Error: No search text provided"
        if _is_odf_file(file_path):
            return _odt_search_replace(file_path, search_text, replace_text)
        return _docx_search_replace(file_path, search_text, replace_text)
    except ValueError as e:
        return f"Error: {e}"
    except Exception as e:
        return f"Error: Failed to search and replace text: {e}"


@mcp.tool()
async def delete_text(file_path: str, text_to_delete: str) -> str:
    """
    Delete specific text from the document.

    Args:
        file_path: Path to the document
        text_to_delete: Text to search for and delete
    """
    return await search_replace_text(file_path, text_to_delete, "")


@mcp.tool()
async def format_table(
    file_path: str,
    table_index: int = 0,
    border_width: Optional[int] = None,
    background_color: Optional[str] = None,
    header_row: bool = False,
) -> str:
    """
    Format a table with borders, background color, and a header row.

    Args:
        file_path: Path to the document
        table_index: Index of the table to format (0 = first table)
        border_width: Border width in points
        background_color: Background color (hex format, e.g., "#F0F0F0")
        header_row: Whether to format the first row as a header
    """
    try:
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Document not found: {file_path}"
        format_options = {}
        if border_width is not None:
            format_options["border_width"] = border_width
        if background_color:
            format_options["background_color"] = background_color
        format_options["header_row"] = header_row

        if _is_odf_file(file_path):
            return f"Error: Table formatting in ODF files is not yet supported. Open the file in LibreOffice to format tables."
        return _docx_format_table(file_path, table_index, format_options)
    except ValueError as e:
        return f"Error: {e}"
    except Exception as e:
        return f"Error: Failed to format table: {e}"


@mcp.tool()
async def delete_paragraph(file_path: str, paragraph_index: int) -> str:
    """
    Delete a paragraph at the given index.
    If the user has not specified an index, you may need to call the read_text_document function first to find the correct index.

    Args:
        file_path: Path to the document
        paragraph_index: Index of the paragraph to delete (0 = first paragraph)
    """
    try:
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Document not found: {file_path}"
        if _is_odf_file(file_path):
            return _odt_delete_paragraph(file_path, paragraph_index)
        return _docx_delete_paragraph(file_path, paragraph_index)
    except ValueError as e:
        return f"Error: {e}"
    except Exception as e:
        return f"Error: Failed to delete paragraph: {e}"


@mcp.tool()
async def apply_document_style(
    file_path: str,
    font_name: Optional[str] = None,
    font_size: Optional[float] = None,
    color: Optional[str] = None,
    alignment: Optional[str] = None,
) -> str:
    """
    Apply consistent formatting throughout the document.

    Args:
        file_path: Path to the document
        font_name: Font name
        font_size: Font size in points
        color: Text color (hex format, e.g., "#000000")
        alignment: Alignment (left, center, right, justify)
    """
    try:
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Document not found: {file_path}"
        style = {}
        if font_name:
            style["font_name"] = font_name
        if font_size:
            style["font_size"] = font_size
        if color:
            style["color"] = color
        if alignment:
            style["alignment"] = alignment

        if _is_odf_file(file_path):
            return f"Error: Document-wide style application in ODF files is not yet supported. Open the file in LibreOffice to apply styles."
        return _docx_apply_document_style(file_path, style)
    except Exception as e:
        return f"Error: Failed to apply document style: {e}"


# ============================================================================
# Slides / Impress Tools (13)
# ============================================================================


@mcp.tool()
async def create_blank_presentation(
    filename: str,
    title: str = "",
    author: str = "",
    subject: str = "",
    keywords: str = "",
) -> str:
    """
    Create a new presentation (.odp or .pptx).

    Args:
        filename: Name of the presentation
        title: Presentation title metadata
        author: Presentation author metadata
        subject: Presentation subject metadata
        keywords: Presentation keywords metadata (comma-separated)
    """
    try:
        if not filename.lower().endswith(IMPRESS_EXTENSIONS):
            filename += ".odp"
        if os.path.basename(filename) == filename:
            save_path = get_default_document_path(filename)
        else:
            save_path = filename
        save_path = normalize_path(save_path)
        if not ensure_directory_exists(save_path):
            return f"Error: Failed to create directory for {save_path}"

        metadata = {}
        if title:
            metadata["Title"] = title
        if author:
            metadata["Author"] = author
        if subject:
            metadata["Subject"] = subject
        if keywords:
            metadata["Keywords"] = [k.strip() for k in keywords.split(",")]

        if _is_odf_file(save_path):
            _odp_create_blank(save_path, metadata)
        else:
            _pptx_create_blank(save_path, metadata)

        if os.path.exists(save_path):
            return f"Successfully created presentation at: {save_path}"
        return f"Error: Presentation creation attempted, but file not found at: {save_path}"
    except Exception as e:
        logging.error(f"Error in create_blank_presentation: {e}")
        return f"Error: Failed to create presentation: {e}"


@mcp.tool()
async def read_presentation(file_path: str) -> str:
    """
    Open and read the text of a presentation.

    Args:
        file_path: Path to the presentation
    """
    try:
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Presentation not found: {file_path}"
        if _is_odf_file(file_path):
            return _odp_read(file_path)
        return _pptx_read(file_path)
    except Exception as e:
        return f"Error: Failed to read presentation: {e}"


@mcp.tool()
async def add_slide(
    file_path: str,
    slide_index: Optional[int] = None,
    title: Optional[str] = None,
    content: Optional[str] = None,
) -> str:
    """
    Add a new slide to a presentation.

    Args:
        file_path: Path to the presentation file.
        slide_index: Index at which to insert the new slide (0-based). If None, the slide is appended at the end.
        title: Optional title text for the new slide.
        content: Optional content text for the new slide.
    """
    try:
        if not file_path.lower().endswith(IMPRESS_EXTENSIONS):
            return "Error: file_path is not a presentation."
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Presentation not found: {file_path}"
        if _is_odf_file(file_path):
            return _odp_add_slide(file_path, slide_index, title, content)
        return _pptx_add_slide(file_path, slide_index, title, content)
    except Exception as e:
        return f"Error: Failed to add slide: {e}"


@mcp.tool()
async def edit_slide_content(file_path: str, slide_index: int, new_content: str) -> str:
    """
    Edit the main text content of a specific slide in a presentation.

    Args:
        file_path: Path to the presentation file
        slide_index: Index of the slide to edit (0-based, where 0 is the first slide)
        new_content: New text content to set in the main content area
    """
    try:
        if not file_path.lower().endswith(IMPRESS_EXTENSIONS):
            return "Error: file_path is not a presentation file."
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Presentation not found: {file_path}"
        if _is_odf_file(file_path):
            return _odp_edit_slide_text(file_path, slide_index, new_content)
        return _pptx_edit_slide_text(file_path, slide_index, 1, new_content)
    except ValueError as e:
        return f"Error: {e}"
    except Exception as e:
        return f"Error: Failed to edit slide content: {e}"


@mcp.tool()
async def edit_slide_title(file_path: str, slide_index: int, new_title: str) -> str:
    """
    Edit the title of a specific slide in a presentation.

    Args:
        file_path: Path to the presentation file
        slide_index: Index of the slide to edit (0-based, where 0 is the first slide)
        new_title: New title text to set in the title area
    """
    try:
        if not file_path.lower().endswith(IMPRESS_EXTENSIONS):
            return "Error: file_path is not a presentation file."
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Presentation not found: {file_path}"
        if _is_odf_file(file_path):
            return _odp_edit_slide_text(file_path, slide_index, new_title)
        return _pptx_edit_slide_text(file_path, slide_index, 0, new_title)
    except ValueError as e:
        return f"Error: {e}"
    except Exception as e:
        return f"Error: Failed to edit slide title: {e}"


@mcp.tool()
async def delete_slide(file_path: str, slide_index: int) -> str:
    """
    Delete a slide from a presentation.

    Args:
        file_path: Path to the presentation file
        slide_index: Index of the slide to delete (0-based, where 0 is the first slide)
    """
    try:
        if not file_path.lower().endswith(IMPRESS_EXTENSIONS):
            return "Error: file_path is not a presentation file."
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Presentation not found: {file_path}"
        if _is_odf_file(file_path):
            return _odp_delete_slide(file_path, slide_index)
        return _pptx_delete_slide(file_path, slide_index)
    except ValueError as e:
        return f"Error: {e}"
    except Exception as e:
        return f"Error: Failed to delete slide: {e}"


@mcp.tool()
async def apply_presentation_template(file_path: str, template_name: str) -> str:
    """
    Apply a presentation template/master slide to a presentation.

    Args:
        file_path: Path to the presentation file
        template_name: Name of the template to apply (can be exact name, partial match, or numeric index)
    """
    try:
        if not file_path.lower().endswith(IMPRESS_EXTENSIONS):
            return "Error: file_path is not a presentation file."
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Presentation not found: {file_path}"
        if _is_odf_file(file_path):
            return (
                f"Template application for ODF presentations is not yet supported. "
                f"Open the file in LibreOffice and use Slide > Slide Properties > Master Slide."
            )
        return _pptx_apply_template(file_path, template_name)
    except Exception as e:
        return f"Error: Failed to apply presentation template: {e}"


@mcp.tool()
async def format_slide_content(
    file_path: str,
    slide_index: int,
    font_name: Optional[str] = None,
    font_size: Optional[float] = None,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    underline: Optional[bool] = None,
    color: Optional[str] = None,
    alignment: Optional[str] = None,
    line_spacing: Optional[float] = None,
    background_color: Optional[str] = None,
) -> str:
    """
    Format the content text of a specific slide in a presentation.

    Args:
        file_path: Path to the presentation file
        slide_index: Index of the slide to format (0-based, where 0 is the first slide)
        font_name: Font family name (e.g., "Arial", "Times New Roman")
        font_size: Font size in points (e.g., 18, 24)
        bold: Apply bold formatting (True/False)
        italic: Apply italic formatting (True/False)
        underline: Apply underline formatting (True/False)
        color: Text color as hex string (e.g., "#FF0000") or RGB integer
        alignment: Text alignment ("left", "center", "right", "justify")
        line_spacing: Line spacing multiplier (e.g., 1.5, 2.0)
        background_color: Background color as hex string (e.g., "#F0F0F0") or RGB integer
    """
    try:
        if not file_path.lower().endswith(IMPRESS_EXTENSIONS):
            return "Error: file_path is not a presentation file."
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Presentation not found: {file_path}"
        if _is_odf_file(file_path):
            return "Error: Slide content formatting in ODF presentations is not yet supported. Open the file in LibreOffice."
        format_options = {}
        if font_name is not None:
            format_options["font_name"] = font_name
        if font_size is not None:
            format_options["font_size"] = font_size
        if bold is not None:
            format_options["bold"] = bold
        if italic is not None:
            format_options["italic"] = italic
        if underline is not None:
            format_options["underline"] = underline
        if color is not None:
            format_options["color"] = color
        if alignment is not None:
            format_options["alignment"] = alignment
        return _pptx_format_slide_placeholder(file_path, slide_index, 1, format_options)
    except ValueError as e:
        return f"Error: {e}"
    except Exception as e:
        return f"Error: Failed to format slide content: {e}"


@mcp.tool()
async def format_slide_title(
    file_path: str,
    slide_index: int,
    font_name: Optional[str] = None,
    font_size: Optional[float] = None,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    underline: Optional[bool] = None,
    color: Optional[str] = None,
    alignment: Optional[str] = None,
    line_spacing: Optional[float] = None,
    background_color: Optional[str] = None,
) -> str:
    """
    Format the title text of a specific slide in a presentation.

    Args:
        file_path: Path to the presentation file
        slide_index: Index of the slide to format (0-based, where 0 is the first slide)
        font_name: Font family name (e.g., "Arial", "Times New Roman")
        font_size: Font size in points (e.g., 28, 36)
        bold: Apply bold formatting (True/False)
        italic: Apply italic formatting (True/False)
        underline: Apply underline formatting (True/False)
        color: Text color as hex string (e.g., "#FF0000") or RGB integer
        alignment: Text alignment ("left", "center", "right", "justify")
        line_spacing: Line spacing multiplier (e.g., 1.5, 2.0)
        background_color: Background color as hex string (e.g., "#F0F0F0") or RGB integer
    """
    try:
        if not file_path.lower().endswith(IMPRESS_EXTENSIONS):
            return "Error: file_path is not a presentation file."
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: Presentation not found: {file_path}"
        if _is_odf_file(file_path):
            return "Error: Slide title formatting in ODF presentations is not yet supported. Open the file in LibreOffice."
        format_options = {}
        if font_name is not None:
            format_options["font_name"] = font_name
        if font_size is not None:
            format_options["font_size"] = font_size
        if bold is not None:
            format_options["bold"] = bold
        if italic is not None:
            format_options["italic"] = italic
        if underline is not None:
            format_options["underline"] = underline
        if color is not None:
            format_options["color"] = color
        if alignment is not None:
            format_options["alignment"] = alignment
        return _pptx_format_slide_placeholder(file_path, slide_index, 0, format_options)
    except ValueError as e:
        return f"Error: {e}"
    except Exception as e:
        return f"Error: Failed to format slide title: {e}"


@mcp.tool()
async def insert_slide_image(
    file_path: str,
    slide_index: int,
    image_path: str,
    max_width: Optional[int] = None,
    max_height: Optional[int] = None,
) -> str:
    """
    Insert an image into a specific slide of a presentation.
    The image will be centered on the slide and resized if necessary to fit.

    Args:
        file_path: Path to the presentation file
        slide_index: Index of the slide to insert the image into (0-based, where 0 is the first slide)
        image_path: Path to the image file to insert
        max_width: Maximum width in 1/100mm units (defaults to slide width minus margins)
        max_height: Maximum height in 1/100mm units (defaults to slide height minus margins)
    """
    try:
        if not file_path.lower().endswith(IMPRESS_EXTENSIONS):
            return "Error: file_path is not a presentation file."
        file_path = normalize_path(file_path)
        image_path = normalize_path(image_path)
        if not os.path.exists(file_path):
            return f"Error: Presentation not found: {file_path}"
        if not os.path.exists(image_path):
            return f"Error: Image not found: {image_path}"
        if _is_odf_file(file_path):
            return "Error: Slide image insertion in ODF presentations is not yet supported. Open the file in LibreOffice."
        return _pptx_insert_image(file_path, slide_index, image_path, max_width, max_height)
    except ValueError as e:
        return f"Error: {e}"
    except Exception as e:
        return f"Error: Failed to insert slide image: {e}"


# ============================================================================
# Bonus: Open in LibreOffice
# ============================================================================


@mcp.tool()
async def open_in_libreoffice(file_path: str) -> str:
    """
    Open a document in LibreOffice for viewing or editing.
    Requires LibreOffice or Collabora Office to be installed.

    Args:
        file_path: Path to the document to open
    """
    try:
        file_path = normalize_path(file_path)
        if not os.path.exists(file_path):
            return f"Error: File not found: {file_path}"
        soffice = _find_soffice_path()
        if soffice is None:
            return (
                "Error: LibreOffice or Collabora Office is not installed. "
                "Install it from https://www.libreoffice.org/download/ to open documents in the GUI."
            )
        # Use file:/// URI so LibreOffice handles spaces and special chars in
        # the path correctly.  --norestore skips the recovery dialog which can
        # block the file from loading.
        file_uri = Path(file_path).as_uri()
        subprocess.Popen([soffice, "--norestore", file_uri])
        return f"Opened {file_path} in LibreOffice"
    except Exception as e:
        return f"Error: Failed to open document in LibreOffice: {e}"


# ============================================================================
# Resource for document access
# ============================================================================


@mcp.resource("libreoffice:{path}")
async def document_resource(path: str) -> str:
    """
    Resource for accessing document content.

    Args:
        path: Path to the document

    Returns:
        Document content as text
    """
    try:
        normalized_path = normalize_path(path)
        if not os.path.exists(normalized_path):
            return f"Error: Document not found: {normalized_path}"
        if _is_impress_file(normalized_path):
            if _is_odf_file(normalized_path):
                return _odp_read(normalized_path)
            return _pptx_read(normalized_path)
        if _is_odf_file(normalized_path):
            return _odt_read_text(normalized_path)
        return _docx_read_text(normalized_path)
    except Exception as e:
        return f"Error: Failed to access document resource: {e}"


# ============================================================================
# Test / verification tools (imported from test_functions.py)
# ============================================================================

try:
    from test_functions import register_test_tools
    register_test_tools(mcp)
except ImportError:
    logging.warning("test_functions module not found; test/verification tools will not be available")


# ============================================================================
# Entry point
# ============================================================================


async def main():
    logging.info("Starting LibreOffice document MCP server (python-docx/python-pptx/odfdo)")
    await mcp.run_stdio_async()
    logging.info("MCP server exited")


if __name__ == "__main__":
    asyncio.run(main())

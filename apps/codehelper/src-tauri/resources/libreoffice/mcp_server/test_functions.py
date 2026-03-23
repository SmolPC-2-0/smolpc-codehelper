"""
Verification / test functions for document inspection.

These are read-only tools exposed via MCP so that the Rust test infrastructure
can verify document operations round-trip correctly.
"""
import json
import logging
import os

from docx import Document as DocxDocument
from docx.oxml.ns import qn

from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _is_odf(path: str) -> bool:
    return path.lower().endswith((".odt", ".odp", ".ods"))


def _is_pptx(path: str) -> bool:
    return path.lower().endswith((".pptx", ".ppsx", ".ppmx", ".potx", ".pomx", ".ppt", ".pps"))


# ---------------------------------------------------------------------------
# 1. get_text_formatting
# ---------------------------------------------------------------------------

def _get_text_formatting_docx(file_path: str, text_to_find: str) -> dict:
    doc = DocxDocument(file_path)
    text_lower = text_to_find.lower()
    for paragraph in doc.paragraphs:
        if text_lower not in paragraph.text.lower():
            continue
        for run in paragraph.runs:
            if text_lower in run.text.lower():
                font = run.font
                color_val = None
                if font.color and font.color.rgb:
                    color_val = str(font.color.rgb)
                alignment = None
                if paragraph.alignment is not None:
                    align_map = {0: "left", 1: "center", 2: "right", 3: "justify"}
                    alignment = align_map.get(int(paragraph.alignment), str(paragraph.alignment))

                return {
                    "found": True,
                    "font_name": font.name,
                    "font_size": float(font.size.pt) if font.size else None,
                    "bold": bool(run.bold) if run.bold is not None else False,
                    "italic": bool(run.italic) if run.italic is not None else False,
                    "underline": bool(run.underline) if run.underline is not None else False,
                    "color": color_val,
                    "alignment": alignment,
                    "paragraph_style": paragraph.style.name if paragraph.style else None,
                    "occurrences_found": paragraph.text.lower().count(text_lower),
                }
    return {"found": False, "error": f"Text '{text_to_find}' not found"}


def get_text_formatting(file_path: str, text_to_find: str) -> str:
    if _is_odf(file_path):
        return json.dumps({"found": False, "error": "ODF text formatting inspection not supported yet"})
    return json.dumps(_get_text_formatting_docx(file_path, text_to_find))


# ---------------------------------------------------------------------------
# 2. get_table_info
# ---------------------------------------------------------------------------

def _get_table_info_docx(file_path: str, table_index: int) -> dict:
    doc = DocxDocument(file_path)
    if table_index >= len(doc.tables):
        return {"error": f"Table index {table_index} out of range (document has {len(doc.tables)} tables)"}
    table = doc.tables[table_index]
    rows = len(table.rows)
    cols = len(table.columns)
    data = []
    for row in table.rows:
        row_data = [cell.text for cell in row.cells]
        data.append(row_data)

    # Try to read border info
    border_info = {}
    tbl = table._tbl
    tblPr = tbl.tblPr
    if tblPr is not None:
        borders_el = tblPr.find(qn("w:tblBorders"))
        if borders_el is not None:
            for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
                el = borders_el.find(qn(f"w:{edge}"))
                if el is not None:
                    border_info[edge] = {
                        "val": el.get(qn("w:val")),
                        "sz": el.get(qn("w:sz")),
                        "color": el.get(qn("w:color")),
                    }

    return {
        "rows": rows,
        "columns": cols,
        "data": data,
        "total_tables": len(doc.tables),
        "borders": border_info,
    }


def get_table_info(file_path: str, table_index: int) -> str:
    if _is_odf(file_path):
        return json.dumps({"error": "ODF table inspection not supported yet"})
    return json.dumps(_get_table_info_docx(file_path, table_index))


# ---------------------------------------------------------------------------
# 3. has_image
# ---------------------------------------------------------------------------

def _has_image_docx(file_path: str) -> dict:
    doc = DocxDocument(file_path)
    image_count = len(doc.inline_shapes)
    result = {
        "has_image": image_count > 0,
        "image_count": image_count,
    }
    if image_count > 0:
        shape = doc.inline_shapes[0]
        result["first_image_width"] = shape.width
        result["first_image_height"] = shape.height
    return result


def has_image(file_path: str) -> str:
    if _is_odf(file_path):
        return json.dumps({"has_image": False, "error": "ODF image inspection not supported yet"})
    return json.dumps(_has_image_docx(file_path))


# ---------------------------------------------------------------------------
# 4. get_page_break_info
# ---------------------------------------------------------------------------

def _get_page_break_info_docx(file_path: str) -> dict:
    doc = DocxDocument(file_path)
    page_breaks = []
    total_breaks = 0

    for i, paragraph in enumerate(doc.paragraphs):
        has_break = False
        # Check for w:br with type="page" in runs
        for run in paragraph.runs:
            for br in run._element.findall(qn("w:br")):
                if br.get(qn("w:type")) == "page":
                    has_break = True
                    total_breaks += 1
        # Check paragraph properties for page break before
        pPr = paragraph._element.find(qn("w:pPr"))
        if pPr is not None:
            # Check sectPr for page break
            for child in pPr:
                if child.tag == qn("w:pageBreakBefore"):
                    has_break = True
                    total_breaks += 1

        if has_break:
            page_breaks.append({
                "paragraph_index": i,
                "has_break": True,
            })

    return {
        "total_page_breaks": total_breaks,
        "paragraph_page_breaks": page_breaks,
        "paragraph_count": len(doc.paragraphs),
    }


def get_page_break_info(file_path: str) -> str:
    if _is_odf(file_path):
        return json.dumps({"error": "ODF page break inspection not supported yet"})
    return json.dumps(_get_page_break_info_docx(file_path))


# ---------------------------------------------------------------------------
# 5. get_presentation_template_info
# ---------------------------------------------------------------------------

def _get_template_info_pptx(file_path: str) -> str:
    prs = PptxPresentation(file_path)
    if prs.slide_masters:
        master = prs.slide_masters[0]
        # Try to get master slide name
        name = master.name if hasattr(master, "name") else None
        if name:
            return f"Template: {name}"
    props = prs.core_properties
    if props.title:
        return f"Template: {props.title}"
    return "Template: Default"


def get_presentation_template_info(file_path: str) -> str:
    if _is_odf(file_path):
        return "Template inspection for ODF not supported yet"
    return _get_template_info_pptx(file_path)


# ---------------------------------------------------------------------------
# 6. get_presentation_text_formatting
# ---------------------------------------------------------------------------

def _get_pptx_text_formatting(file_path: str, text_to_find: str, slide_index=None) -> dict:
    prs = PptxPresentation(file_path)
    text_lower = text_to_find.lower()

    slides_to_check = []
    if slide_index is not None:
        if slide_index < len(prs.slides):
            slides_to_check = [(slide_index, prs.slides[slide_index])]
    else:
        slides_to_check = list(enumerate(prs.slides))

    for idx, slide in slides_to_check:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if text_lower in run.text.lower():
                        font = run.font
                        color_val = None
                        if font.color and font.color.rgb:
                            color_val = str(font.color.rgb)
                        alignment = None
                        if paragraph.alignment is not None:
                            align_map = {0: "left", 1: "center", 2: "right", 3: "justify"}
                            alignment = align_map.get(int(paragraph.alignment), str(paragraph.alignment))

                        # Determine shape category
                        shape_cat = "unknown"
                        if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                            ph_idx = shape.placeholder_format.idx
                            shape_cat = "title" if ph_idx == 0 else "content"

                        return {
                            "found": True,
                            "found_on_slide": idx,
                            "shape_category": shape_cat,
                            "font_name": font.name,
                            "font_size": float(font.size.pt) if font.size else None,
                            "bold": bool(font.bold) if font.bold is not None else False,
                            "italic": bool(font.italic) if font.italic is not None else False,
                            "underline": bool(font.underline) if font.underline is not None else False,
                            "color": color_val,
                            "alignment": alignment,
                        }

    return {"found": False, "error": f"Text '{text_to_find}' not found"}


def get_presentation_text_formatting(file_path: str, text_to_find: str, slide_index=None) -> str:
    if _is_odf(file_path):
        return json.dumps({"found": False, "error": "ODF presentation text formatting not supported yet"})
    return json.dumps(_get_pptx_text_formatting(file_path, text_to_find, slide_index))


# ---------------------------------------------------------------------------
# 7. get_slide_image_info
# ---------------------------------------------------------------------------

def _get_slide_image_info_pptx(file_path: str, slide_index: int) -> dict:
    prs = PptxPresentation(file_path)
    if slide_index >= len(prs.slides):
        return {"error": f"Slide index {slide_index} out of range"}
    slide = prs.slides[slide_index]
    images = []
    for i, shape in enumerate(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            images.append({
                "shape_index": i,
                "width": shape.width,
                "height": shape.height,
                "left": shape.left,
                "top": shape.top,
                "width_mm": round(shape.width / 36000, 2),
                "height_mm": round(shape.height / 36000, 2),
            })

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    return {
        "slide_index": slide_index,
        "has_images": len(images) > 0,
        "image_count": len(images),
        "images": images,
        "total_slides": len(prs.slides),
        "slide_width": slide_width,
        "slide_height": slide_height,
    }


def get_slide_image_info(file_path: str, slide_index: int) -> str:
    if _is_odf(file_path):
        return json.dumps({"error": "ODF slide image inspection not supported yet"})
    return json.dumps(_get_slide_image_info_pptx(file_path, slide_index))


# ---------------------------------------------------------------------------
# Registration — called from mcp_server.py
# ---------------------------------------------------------------------------

def register_test_tools(mcp_server):
    """Register all test/verification functions as MCP tools."""

    @mcp_server.tool()
    async def test_get_text_formatting(file_path: str, text_to_find: str) -> str:
        """Get text formatting details for test verification."""
        return get_text_formatting(file_path, text_to_find)

    @mcp_server.tool()
    async def test_get_table_info(file_path: str, table_index: int = 0) -> str:
        """Get table structure and formatting for test verification."""
        return get_table_info(file_path, table_index)

    @mcp_server.tool()
    async def test_has_image(file_path: str) -> str:
        """Check if document contains images for test verification."""
        return has_image(file_path)

    @mcp_server.tool()
    async def test_get_page_break_info(file_path: str) -> str:
        """Get page break details for test verification."""
        return get_page_break_info(file_path)

    @mcp_server.tool()
    async def test_get_presentation_template_info(file_path: str) -> str:
        """Get presentation template info for test verification."""
        return get_presentation_template_info(file_path)

    @mcp_server.tool()
    async def test_get_presentation_text_formatting(
        file_path: str, text_to_find: str, slide_index: int = -1
    ) -> str:
        """Get presentation text formatting for test verification."""
        si = slide_index if slide_index >= 0 else None
        return get_presentation_text_formatting(file_path, text_to_find, si)

    @mcp_server.tool()
    async def test_get_slide_image_info(file_path: str, slide_index: int) -> str:
        """Get slide image info for test verification."""
        return get_slide_image_info(file_path, slide_index)
